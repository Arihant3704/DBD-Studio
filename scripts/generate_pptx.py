import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Color Constants for Light Theme
COLOR_BG = RGBColor(249, 250, 251)
COLOR_TEXT = RGBColor(31, 41, 55)
COLOR_AMBER = RGBColor(217, 119, 6) # Using slightly darker amber for text readability
COLOR_LIGHT_AMBER = RGBColor(245, 158, 11)

def add_title(slide, text, left, top, width, height, font_size=44, color=COLOR_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = 'Georgia'
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color
    return txBox

def add_subtitle(slide, text, left, top, width, height, font_size=18, color=RGBColor(75, 85, 99)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.name = 'Arial'
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    return txBox

def set_slide_bg(slide, color=COLOR_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

blank_slide_layout = prs.slide_layouts[6]

# Slide 1: Cover
slide1 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide1, color=RGBColor(255, 255, 255))

if os.path.exists("photo-video/WhatsApp Image 2025-08-11 at 10.51.15 (1).jpeg"):
    slide1.shapes.add_picture("photo-video/WhatsApp Image 2025-08-11 at 10.51.15 (1).jpeg", Inches(8), Inches(0), height=Inches(9))

add_title(slide1, "DBD Studio:\nArchitectural Curation", Inches(1), Inches(2), Inches(7), Inches(3), font_size=64)
add_subtitle(slide1, "Transforming brand identity into immersive physical narratives.", Inches(1), Inches(5.5), Inches(6), Inches(1))

# Slide 2: The Visionary
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide2)
add_title(slide2, "The Visionary.", Inches(8), Inches(2), Inches(7), Inches(1))
add_subtitle(slide2, "Beulah Hephzibah\nRedefining space as an experience through atmospheric minimalism.", Inches(8), Inches(3.5), Inches(7), Inches(2))
if os.path.exists("photo-video/new/brown_boy_cafe.jpeg"):
    slide2.shapes.add_picture("photo-video/new/brown_boy_cafe.jpeg", Inches(0), Inches(0), width=Inches(7.5), height=Inches(9))

# Slide 3: Retail Portfolio
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide3, color=RGBColor(255, 255, 255))
add_title(slide3, "Retail Narrative & Spatial Branding", Inches(1), Inches(0.5), Inches(14), Inches(1))
if os.path.exists("photo-video/CHUMBAK.jpeg"):
    slide3.shapes.add_picture("photo-video/CHUMBAK.jpeg", Inches(1), Inches(2), width=Inches(6.5), height=Inches(5))
if os.path.exists("photo-video/unused/WhatsApp Image 2025-08-11 at 10.51.32 (2).jpeg"):
    slide3.shapes.add_picture("photo-video/unused/WhatsApp Image 2025-08-11 at 10.51.32 (2).jpeg", Inches(8.5), Inches(2), width=Inches(6.5), height=Inches(5))

# Slide 4: Culinary Hubs
slide4 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide4)
add_title(slide4, "Culinary Hubs.", Inches(1), Inches(1), Inches(10), Inches(1))
if os.path.exists("photo-video/unused/brown_boy_cafe_interior.jpeg"):
    slide4.shapes.add_picture("photo-video/unused/brown_boy_cafe_interior.jpeg", Inches(1), Inches(2.2), width=Inches(4.5), height=Inches(5))
if os.path.exists("photo-video/new/sandowitch_inside_view.jpeg"):
    slide4.shapes.add_picture("photo-video/new/sandowitch_inside_view.jpeg", Inches(6), Inches(2.2), width=Inches(4.5), height=Inches(5))
add_subtitle(slide4, "F&B Strategy:\nPartnering with gourmet hubs to define atmospheric dining.", Inches(11), Inches(2.2), Inches(4), Inches(2), color=COLOR_AMBER)

# Slide 5: Luxury Residential
slide5 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide5, color=RGBColor(255, 255, 255))
add_title(slide5, "Luxury Residential.", Inches(1), Inches(0.8), Inches(10), Inches(1))
img_configs = [
    ("photo-video/new/the_trost_inside.jpeg", 0.5),
    ("photo-video/new/the_trost_outside.jpeg", 4.3),
    ("photo-video/unused/sl_shet_interior.png", 8.1),
    ("photo-video/SIGNATURE FELICITY.jpeg", 11.9)
]
for path, left in img_configs:
    if os.path.exists(path):
        slide5.shapes.add_picture(path, Inches(left), Inches(2.5), width=Inches(3.6), height=Inches(4.5))

# Slide 6: Functional Narratives
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide6)
add_title(slide6, "Functional Narratives.", Inches(1), Inches(1), Inches(10), Inches(1))
if os.path.exists("photo-video/DHOBILITE.jpeg"):
    slide6.shapes.add_picture("photo-video/DHOBILITE.jpeg", Inches(1), Inches(2.5), width=Inches(6.5), height=Inches(5))
if os.path.exists("photo-video/new/sandowitch_outsideview.jpeg"):
    slide6.shapes.add_picture("photo-video/new/sandowitch_outsideview.jpeg", Inches(8.5), Inches(2.5), width=Inches(6.5), height=Inches(5))

# Slide 7: The Process
slide7 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide7, color=RGBColor(243, 244, 246))
add_title(slide7, "The Turnkey Pillar", Inches(1), Inches(1), Inches(14), Inches(1))
add_subtitle(slide7, "01. Strategic Analysis  |  02. Creative Curation  |  03. Precision Execution", Inches(1), Inches(3), Inches(14), Inches(1), font_size=28, color=COLOR_AMBER)

# Slide 8: Why DBD (Edge)
slide8 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide8, color=RGBColor(255, 255, 255))
add_title(slide8, "The DBD Edge.", Inches(1), Inches(1), Inches(7), Inches(1))
add_subtitle(slide8, "- ROI Focused Design\n- Sustainable Brutalism\n- Turnkey civil mastery", Inches(1), Inches(2.5), Inches(6), Inches(3))
if os.path.exists("photo-video/unused/dbd_edge_materials.png"):
    slide8.shapes.add_picture("photo-video/unused/dbd_edge_materials.png", Inches(8), Inches(0), width=Inches(8), height=Inches(9))

# Slide 9: Global Vision
slide9 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide9)
add_title(slide9, "Ready to curate your vision?", Inches(1), Inches(3), Inches(14), Inches(2))
add_subtitle(slide9, "Bangalore | NYC | London | Tokyo\nhello@dbdstudio.com", Inches(1), Inches(5), Inches(14), Inches(1), color=COLOR_AMBER)

# Slide 10: Contact Details
slide10 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide10, color=RGBColor(255, 255, 255))
add_title(slide10, "Contact Details.", Inches(1), Inches(1), Inches(10), Inches(1))
add_subtitle(slide10, "Address", Inches(1), Inches(3), Inches(4), Inches(0.5), font_size=12, color=RGBColor(156, 163, 175))
add_subtitle(slide10, "No. 15/1, GK Layout, 4th Cross,\nElectronic City, Bangalore – 560100", Inches(1), Inches(3.6), Inches(4), Inches(2), font_size=18, color=COLOR_TEXT)

add_subtitle(slide10, "Communication", Inches(6), Inches(3), Inches(4), Inches(0.5), font_size=12, color=RGBColor(156, 163, 175))
add_subtitle(slide10, "+91 63617 36075\ndreambuilddesign@outlook.com", Inches(6), Inches(3.6), Inches(4), Inches(2), font_size=18, color=COLOR_TEXT)

add_subtitle(slide10, "Working Hours", Inches(11), Inches(3), Inches(4), Inches(0.5), font_size=12, color=RGBColor(156, 163, 175))
add_subtitle(slide10, "Mon - Fri: 9:00 AM - 6:00 PM\nSat: 10:00 AM - 4:00 PM", Inches(11), Inches(3.6), Inches(4), Inches(2), font_size=18, color=COLOR_TEXT)

# Save
prs.save('company-profile.pptx')
print("PowerPoint generated: company-profile.pptx")
